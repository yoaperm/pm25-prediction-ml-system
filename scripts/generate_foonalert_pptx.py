#!/usr/bin/env python3
"""
generate_foonalert_pptx.py
==========================
Generate FoonAlert presentation (16 slides, 10-min video).
Includes real Airflow results, paper references, and model analysis.

Usage:
    python scripts/generate_foonalert_pptx.py
    -> outputs: reports/FoonAlert_Presentation.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# -- Theme colors --
PRIMARY = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0xE7, 0x4C, 0x3C)
ACCENT2 = RGBColor(0x9B, 0x59, 0xB6)
SUCCESS = RGBColor(0x27, 0xAE, 0x60)
WARNING = RGBColor(0xF3, 0x9C, 0x12)
TEXT_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x2C, 0x3E, 0x50)
BG_LIGHT = RGBColor(0xFA, 0xFA, 0xFA)
GRAY = RGBColor(0x95, 0xA5, 0xA6)
BLUE = RGBColor(0x34, 0x98, 0xDB)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height, *,
             font_size=18, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT,
             font_name="Calibri"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return tb


def add_bullets(slide, bullets, left, top, width, height, *,
                font_size=18, color=TEXT_DARK, font_name="Calibri"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"\u2022  {b}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = font_name
    return tb


def add_box(slide, left, top, width, height, *, fill_color=BG_LIGHT, line_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


# =====================================================================
# SLIDES
# =====================================================================

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY)
    add_text(slide, "FoonAlert", 0.5, 2.5, 12.5, 1.0,
             font_size=60, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    add_text(slide, "Real-Time PM2.5 Spike Forecasting System", 0.5, 3.6, 12.5, 0.6,
             font_size=24, color=RGBColor(0xBD, 0xC3, 0xC7), align=PP_ALIGN.CENTER)
    add_text(slide, "7 Models | Hourly Predictions | Auto-Retrain Pipeline",
             0.5, 4.3, 12.5, 0.6,
             font_size=18, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, "Team: YG  |  Music  |  Sunta  |  Olf  |  Perm",
             0.5, 6.0, 12.5, 0.5,
             font_size=14, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, '"Don\'t just see what happened \u2014 know what\'s about to happen."',
             0.5, 6.5, 12.5, 0.5,
             font_size=16, color=WARNING, align=PP_ALIGN.CENTER)


def slide_hook(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_LIGHT)
    add_text(slide, "The Problem", 0.5, 0.5, 12.5, 0.6,
             font_size=20, color=GRAY)
    add_text(slide, "PM2.5 spike 18 \u2192 108 \u00b5g/m\u00b3 in 8 hours",
             0.5, 1.8, 12.5, 1.0,
             font_size=40, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, "No one warned 11 million people in Bangkok.",
             0.5, 3.0, 12.5, 0.8,
             font_size=28, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, "Current apps: show NOW only | No prediction | No early warning",
             0.5, 4.5, 12.5, 0.6,
             font_size=20, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    add_box(slide, 2.0, 5.3, 9.5, 1.2, fill_color=SUCCESS)
    add_text(slide, "FoonAlert: Predict +1h, +6h, +24h ahead | 7 models competing | Auto-retrain",
             2.0, 5.5, 9.5, 0.8,
             font_size=18, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    add_text(slide, "[ YG: 0:00-1:00 ]", 0.5, 6.9, 12.5, 0.3,
             font_size=11, color=GRAY, align=PP_ALIGN.CENTER)


def slide_research(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_LIGHT)
    add_text(slide, "Research Foundation", 0.5, 0.4, 12.5, 0.5,
             font_size=20, color=GRAY)
    add_text(slide, "3 Papers That Shaped Our Approach",
             0.5, 0.9, 12.5, 0.7, font_size=30, bold=True, color=TEXT_DARK)

    papers = [
        ("[1] Malakouti (2025)",
         "From accurate to actionable: PM2.5 forecasting with feature engineering & SHAP",
         "Environmental Challenges 21, 101290",
         "We used: Lag/rolling/diff features + SHAP interpretation",
         BLUE),
        ("[2] Buya et al. (2024)",
         "Estimating Ground-level Hourly PM2.5 in Thailand (IEEE JSTARS)",
         "DOI: 10.1109/JSTARS.2024.3384964",
         "We used: Hourly granularity + rush-hour time features",
         ACCENT),
        ("[3] Jankondee et al. (2024)",
         "PM2.5 modeling based on CALIPSO in Bangkok",
         "Creative Science 16(3), DOI: 10.55674/cs.v16i3.257117",
         "We used: Linear models as strong baseline for Bangkok PM2.5",
         SUCCESS),
    ]

    y = 1.9
    for title, desc, journal, insight, color in papers:
        add_box(slide, 0.5, y, 12.0, 1.5, fill_color=RGBColor(0xFF, 0xFF, 0xFF), line_color=color)
        add_text(slide, title, 0.7, y + 0.1, 7.5, 0.4,
                 font_size=15, bold=True, color=color)
        add_text(slide, desc, 0.7, y + 0.5, 7.5, 0.5,
                 font_size=13, color=TEXT_DARK)
        add_text(slide, journal, 0.7, y + 1.0, 7.5, 0.3,
                 font_size=11, color=GRAY)
        add_box(slide, 8.5, y + 0.2, 3.8, 1.0, fill_color=color)
        add_text(slide, insight, 8.6, y + 0.35, 3.6, 0.8,
                 font_size=11, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
        y += 1.7

    add_text(slide, "[ Music: 1:00-3:00 ]", 0.5, 6.9, 12.5, 0.3,
             font_size=11, color=GRAY, align=PP_ALIGN.CENTER)


def slide_data(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_LIGHT)
    add_text(slide, "Data Pipeline", 0.5, 0.4, 12.5, 0.5,
             font_size=20, color=GRAY)
    add_text(slide, "AirBKK API \u2192 Airflow \u2192 PostgreSQL \u2192 Triton \u2192 Dashboard",
             0.5, 1.0, 12.5, 0.7, font_size=26, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    boxes = [
        ("AirBKK API", "Hourly\nThai Gov", 0.3, BLUE),
        ("Airflow", "Ingest DAG\nHourly cron", 2.9, WARNING),
        ("PostgreSQL", "5 stations\n96k+ rows", 5.5, PRIMARY),
        ("Training", "7 models\nAuto-retrain", 8.1, ACCENT),
        ("Triton+API", "ONNX serve\n<10ms", 10.7, ACCENT2),
    ]
    for label, desc, x, color in boxes:
        add_box(slide, x, 2.2, 2.3, 1.6, fill_color=color)
        add_text(slide, label, x, 2.3, 2.3, 0.5,
                 font_size=14, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x, 2.8, 2.3, 0.8,
                 font_size=11, color=RGBColor(0xDD, 0xDD, 0xFF), align=PP_ALIGN.CENTER)

    add_text(slide, "Stations: 56, 57, 58, 59, 61  |  Jan 2023 \u2192 May 2026  |  Hourly  |  96,000+ rows",
             0.5, 4.2, 12.5, 0.5, font_size=16, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    add_box(slide, 0.5, 5.0, 12.0, 1.7, fill_color=RGBColor(0xFF, 0xFA, 0xE5))
    add_text(slide, "Design Decisions (from papers)", 0.7, 5.1, 11.6, 0.4,
             font_size=14, bold=True, color=WARNING)
    add_bullets(slide, [
        "Hourly (not daily): PM2.5 spike occurs in 2-4h [Buya 2024]",
        "Feature eng: lag 1-24h + rolling mean/std + time [Malakouti 2025]",
        "Linear baseline strong in Bangkok: R\u00b2=0.99 [Jankondee 2024]",
    ], 0.7, 5.5, 11.6, 1.2, font_size=14, color=TEXT_DARK)

    add_text(slide, "[ Music: 1:00-3:00 ]", 0.5, 6.9, 12.5, 0.3,
             font_size=11, color=GRAY, align=PP_ALIGN.CENTER)


def slide_models(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_LIGHT)
    add_text(slide, "7 Models \u2014 The Battle Lineup", 0.5, 0.4, 12.5, 0.7,
             font_size=30, bold=True, color=TEXT_DARK)

    # Row 1: ML models
    row1 = [
        ("Linear", "Baseline\nRMSE 8.51", GRAY, "Sunta"),
        ("Ridge", "BEST\nRMSE 8.50", SUCCESS, "Sunta"),
        ("Random Forest", "Ensemble\nRMSE 8.81", BLUE, "Sunta"),
        ("XGBoost", "Boosting\nRMSE 8.64", WARNING, "Sunta"),
    ]
    for i, (name, desc, color, owner) in enumerate(row1):
        x = 0.3 + i * 3.2
        add_box(slide, x, 1.4, 3.0, 1.8, fill_color=RGBColor(0xFF, 0xFF, 0xFF), line_color=color)
        add_text(slide, name, x, 1.5, 3.0, 0.4,
                 font_size=15, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x, 1.9, 3.0, 0.8,
                 font_size=13, color=TEXT_DARK, align=PP_ALIGN.CENTER)
        add_text(slide, owner, x, 2.8, 3.0, 0.3,
                 font_size=10, color=GRAY, align=PP_ALIGN.CENTER)

    # Row 2: Deep/Statistical
    row2 = [
        ("LSTM", "Deep Learning\nRMSE 10.53\n(needs more data)", ACCENT, "Sunta"),
        ("SARIMA", "Statistical\nRMSE 8.90\n(spike timing)", BLUE, "Olf"),
        ("Transformer", "Self-Attention\nRMSE 8.82\n(early detection)", ACCENT2, "Perm"),
    ]
    for i, (name, desc, color, owner) in enumerate(row2):
        x = 1.2 + i * 3.8
        add_box(slide, x, 3.7, 3.5, 2.5, fill_color=RGBColor(0xFF, 0xFF, 0xFF), line_color=color)
        add_text(slide, name, x, 3.8, 3.5, 0.5,
                 font_size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x, 4.3, 3.5, 1.3,
                 font_size=13, color=TEXT_DARK, align=PP_ALIGN.CENTER)
        add_text(slide, owner, x, 5.8, 3.5, 0.3,
                 font_size=11, color=GRAY, align=PP_ALIGN.CENTER)

    add_text(slide, "[ Sunta: 3:00-4:30 | Olf: 4:30-5:30 | Perm: 5:30-6:00 ]",
             0.5, 6.9, 12.5, 0.3, font_size=11, color=GRAY, align=PP_ALIGN.CENTER)


def slide_features(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_LIGHT)
    add_text(slide, "Feature Engineering \u2014 19 Features", 0.5, 0.4, 12.5, 0.7,
             font_size=28, bold=True, color=TEXT_DARK)

    cols = [
        ("Lag Features (6)", [
            "pm25_lag_1h", "pm25_lag_2h",
            "pm25_lag_3h", "pm25_lag_6h",
            "pm25_lag_12h", "pm25_lag_24h",
        ]),
        ("Rolling Stats (6)", [
            "rolling_mean_6h", "rolling_mean_12h",
            "rolling_mean_24h", "rolling_std_6h",
            "rolling_std_12h", "rolling_std_24h",
        ]),
        ("Time + Diff (7)", [
            "hour", "day_of_week", "month",
            "day_of_year", "is_weekend",
            "pm25_diff_1h", "pm25_diff_24h",
        ]),
    ]

    for i, (title, feats) in enumerate(cols):
        x = 0.5 + i * 4.2
        add_box(slide, x, 1.5, 4.0, 4.8, fill_color=RGBColor(0xFF, 0xFF, 0xFF), line_color=PRIMARY)
        add_text(slide, title, x, 1.6, 4.0, 0.5,
                 font_size=16, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
        add_bullets(slide, feats, x + 0.2, 2.2, 3.7, 4.0,
                    font_size=13, color=TEXT_DARK, font_name="Consolas")

    add_text(slide, "Target: PM2.5 at T+24h  |  All features shifted to prevent leakage",
             0.5, 6.5, 12.5, 0.4, font_size=14, bold=True,
             color=ACCENT2, align=PP_ALIGN.CENTER)


def slide_demo_marker(prs, title, subtitle, url, scene):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, ACCENT)
    add_text(slide, "LIVE DEMO", 0.5, 1.0, 12.5, 0.8,
             font_size=32, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    add_text(slide, title, 0.5, 2.2, 12.5, 1.0,
             font_size=44, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    add_text(slide, subtitle, 0.5, 3.5, 12.5, 0.8,
             font_size=20, color=RGBColor(0xFF, 0xE0, 0xE0), align=PP_ALIGN.CENTER)
    add_box(slide, 2.5, 4.7, 8.5, 1.5, fill_color=PRIMARY)
    add_text(slide, url, 2.5, 5.0, 8.5, 0.5,
             font_size=18, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    add_text(slide, scene, 2.5, 5.5, 8.5, 0.5,
             font_size=14, color=RGBColor(0xBD, 0xC3, 0xC7), align=PP_ALIGN.CENTER)


def slide_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_LIGHT)
    add_text(slide, "Results \u2014 Real Training (Airflow, Station 56)", 0.5, 0.3, 12.5, 0.7,
             font_size=28, bold=True, color=TEXT_DARK)
    add_text(slide, "T+24h forecast | Test: 3 months | Trained 2026-05-07",
             0.5, 0.9, 12.5, 0.4, font_size=13, color=GRAY)

    headers = ["#", "Model", "RMSE", "MAE", "R\u00b2", "Owner"]
    col_w = [0.6, 3.2, 1.5, 1.5, 1.5, 2.0]
    x0 = 0.7
    y = 1.5
    add_box(slide, x0, y, sum(col_w), 0.55, fill_color=PRIMARY)
    cx = x0
    for h, w in zip(headers, col_w):
        add_text(slide, h, cx, y + 0.08, w, 0.4,
                 font_size=13, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
        cx += w

    rows = [
        ("1", "Ridge Regression", "8.50", "6.39", "0.22", "Sunta"),
        ("2", "Linear Regression", "8.51", "6.39", "0.22", "Sunta"),
        ("3", "XGBoost", "8.64", "6.42", "0.20", "Sunta"),
        ("4", "Random Forest", "8.81", "6.61", "0.17", "Sunta"),
        ("5", "Transformer", "8.82", "6.59", "0.16", "Perm"),
        ("6", "SARIMA", "8.90", "6.52", "0.15", "Olf"),
        ("7", "LSTM", "10.53", "7.55", "-0.19", "Sunta"),
    ]

    for ri, row in enumerate(rows):
        ry = y + 0.55 + ri * 0.5
        bg = RGBColor(0xE8, 0xFF, 0xE8) if ri == 0 else (
             RGBColor(0xF8, 0xF8, 0xF8) if ri % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF))
        add_box(slide, x0, ry, sum(col_w), 0.5, fill_color=bg)
        cx = x0
        for val, w in zip(row, col_w):
            bold = (ri == 0)
            c = SUCCESS if ri == 0 else TEXT_DARK
            add_text(slide, val, cx, ry + 0.08, w, 0.35,
                     font_size=12, bold=bold, color=c, align=PP_ALIGN.CENTER)
            cx += w

    # Summary banner
    add_box(slide, 0.7, 5.3, 11.5, 1.3, fill_color=SUCCESS)
    add_text(slide, "Key Findings:", 0.9, 5.4, 11.0, 0.4,
             font_size=16, bold=True, color=TEXT_LIGHT)
    add_bullets(slide, [
        "Ridge wins overall (simple + good features = best accuracy)",
        "Transformer best for spike early detection (attention mechanism)",
        "LSTM needs more data; simple models beat it here",
    ], 0.9, 5.8, 11.0, 0.8, font_size=13, color=TEXT_LIGHT)

    add_text(slide, "[ YG: 6:30-8:30 ]", 0.5, 6.9, 12.5, 0.3,
             font_size=11, color=GRAY, align=PP_ALIGN.CENTER)


def slide_when_to_use(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_LIGHT)
    add_text(slide, "When to Use Which Model?", 0.5, 0.4, 12.5, 0.7,
             font_size=28, bold=True, color=TEXT_DARK)

    rows = [
        ("Routine monitoring", "Ridge", "Best accuracy, fast inference", SUCCESS),
        ("Spike early detection", "Transformer", "Attention captures onset patterns", ACCENT2),
        ("Seasonal analysis", "SARIMA", "Interpretable seasonal decomposition", BLUE),
        ("Resource-limited", "Ridge / Linear", "Minimal compute, ONNX <2ms", GRAY),
        ("Critical alert", "Ensemble (all)", "Consensus vote = high confidence", ACCENT),
    ]
    y = 1.4
    for scenario, model, why, color in rows:
        add_box(slide, 0.5, y, 12.0, 0.9, fill_color=RGBColor(0xFF, 0xFF, 0xFF), line_color=color)
        add_text(slide, scenario, 0.7, y + 0.2, 4.0, 0.5, font_size=15, color=TEXT_DARK)
        add_text(slide, model, 4.8, y + 0.2, 3.0, 0.5, font_size=15, bold=True, color=color)
        add_text(slide, why, 7.9, y + 0.2, 4.5, 0.5, font_size=13, color=GRAY)
        y += 1.05


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_LIGHT)
    add_text(slide, "Production Architecture", 0.5, 0.4, 12.5, 0.7,
             font_size=28, bold=True, color=TEXT_DARK)

    layers = [
        ("Data Layer", "AirBKK API \u2192 Airflow Hourly Ingest \u2192 PostgreSQL (pm25_raw_hourly)", BLUE),
        ("ML Layer", "Airflow Training DAG \u2192 7 Models \u2192 MLflow tracking \u2192 ONNX export", ACCENT),
        ("Serving Layer", "Triton Inference Server \u2192 FastAPI \u2192 Streamlit FoonAlert UI", ACCENT2),
        ("Monitoring Layer", "PSI drift detection \u2192 Auto-retrain trigger \u2192 Hot-swap via Triton", SUCCESS),
    ]
    y = 1.4
    for title, desc, color in layers:
        add_box(slide, 0.5, y, 12.0, 1.1, fill_color=RGBColor(0xFF, 0xFF, 0xFF), line_color=color)
        add_text(slide, title, 0.7, y + 0.1, 3.0, 0.4, font_size=14, bold=True, color=color)
        add_text(slide, desc, 3.8, y + 0.3, 8.5, 0.5, font_size=14, color=TEXT_DARK)
        y += 1.25

    add_box(slide, 1.5, 6.0, 10.0, 0.7, fill_color=PRIMARY)
    add_text(slide, "Deploy: docker compose up | Zero downtime | Self-healing pipeline",
             1.5, 6.1, 10.0, 0.5, font_size=16, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    add_text(slide, "[ Music: 8:30-9:30 ]", 0.5, 6.9, 12.5, 0.3,
             font_size=11, color=GRAY, align=PP_ALIGN.CENTER)


def slide_why_matters(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY)
    add_text(slide, "Why It Matters", 0.5, 0.5, 12.5, 0.7,
             font_size=20, color=GRAY)
    add_text(slide, "Bangkok needs early warning", 0.5, 1.3, 12.5, 0.8,
             font_size=38, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    stats = [
        ("11M", "people at risk", ACCENT),
        ("100+", "unhealthy days/yr", WARNING),
        ("24h", "advance prediction", SUCCESS),
        ("7", "competing models", ACCENT2),
    ]
    for i, (num, label, color) in enumerate(stats):
        x = 0.5 + i * 3.15
        add_box(slide, x, 3.0, 2.95, 2.5, fill_color=RGBColor(0x24, 0x24, 0x44))
        add_text(slide, num, x, 3.2, 2.95, 1.0,
                 font_size=44, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, label, x, 4.4, 2.95, 0.8,
                 font_size=14, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    add_text(slide, "Early warning \u2192 behavior change \u2192 reduced exposure \u2192 lower health cost",
             0.5, 6.0, 12.5, 0.5, font_size=18, color=RGBColor(0xBD, 0xC3, 0xC7), align=PP_ALIGN.CENTER)


def slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY)
    add_text(slide, "Most apps tell you:", 0.5, 1.5, 12.5, 0.6,
             font_size=22, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, '"How bad is the air NOW?"', 0.5, 2.2, 12.5, 0.8,
             font_size=36, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    add_text(slide, "FoonAlert asks:", 0.5, 3.8, 12.5, 0.6,
             font_size=22, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, '"How bad WILL it become \u2014', 0.5, 4.5, 12.5, 0.8,
             font_size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, 'and can we warn people in time?"', 0.5, 5.2, 12.5, 0.8,
             font_size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, "Thank you!  |  Questions?", 0.5, 6.5, 12.5, 0.5,
             font_size=22, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    add_text(slide, "[ YG: 9:30-10:00 ]", 0.5, 6.9, 12.5, 0.3,
             font_size=11, color=GRAY, align=PP_ALIGN.CENTER)


def slide_references(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_LIGHT)
    add_text(slide, "References", 0.5, 0.4, 12.5, 0.6,
             font_size=28, bold=True, color=TEXT_DARK)

    refs = [
        "[1] Malakouti, S.M. (2025). From accurate to actionable: Interpretable PM2.5",
        "    forecasting with feature engineering and SHAP. Env. Challenges, 21, 101290.",
        "",
        "[2] Buya, S., Gokon, H., Dam, H.C., Usanavasin, S., Karnjana, J. (2024).",
        "    Estimating Ground-level Hourly PM2.5 in Thailand using Satellite Data.",
        "    IEEE J. Sel. Top. Appl. Earth Obs. Remote Sens. DOI:10.1109/JSTARS.2024.3384964",
        "",
        "[3] Jankondee, Y., Kumharn, W., et al. (2024). PM2.5 modeling based on CALIPSO",
        "    in Bangkok. Creative Science, 16(3). DOI:10.55674/cs.v16i3.257117",
    ]
    y = 1.2
    for line in refs:
        if line:
            add_text(slide, line, 0.7, y, 11.5, 0.35,
                     font_size=13, color=TEXT_DARK, font_name="Consolas")
        y += 0.35

    add_text(slide, "Tech Stack", 0.7, 4.7, 11.5, 0.5,
             font_size=16, bold=True, color=PRIMARY)
    add_bullets(slide, [
        "Apache Airflow | Triton Inference Server | ONNX Runtime",
        "PostgreSQL | FastAPI | Streamlit | MLflow",
        "Docker Compose | AirBKK API (Thai air quality)",
    ], 0.7, 5.2, 11.5, 1.5, font_size=14, color=TEXT_DARK)


def slide_qa(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_LIGHT)
    add_text(slide, "Q&A \u2014 Prepared Answers", 0.5, 0.4, 12.5, 0.6,
             font_size=26, bold=True, color=TEXT_DARK)

    qas = [
        ("Q: Why does Ridge beat Deep Learning?",
         "Good features + simple model > complex model. [Malakouti 2025] shows SHAP lag-1 dominates."),
        ("Q: Which model for production?",
         "Ridge (best accuracy). Triton can hot-swap if Transformer/SARIMA improve with more data."),
        ("Q: How often does it retrain?",
         "Daily check. If PSI > 0.2 or MAE > threshold -> auto-trigger training DAG."),
        ("Q: Why LSTM performs poorly?",
         "Needs more data + longer sequences. With 5+ years data, it should improve significantly."),
        ("Q: What's novel here?",
         "7 models competing on same pipeline + auto-monitoring + production deployment (not just research)."),
    ]
    y = 1.2
    for q, a in qas:
        add_text(slide, q, 0.5, y, 12.0, 0.35, font_size=14, bold=True, color=PRIMARY)
        add_text(slide, a, 0.7, y + 0.35, 11.8, 0.45, font_size=13, color=TEXT_DARK)
        y += 0.95


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 16 slides for 10-min video
    slide_title(prs)          # 1
    slide_hook(prs)           # 2
    slide_research(prs)       # 3 - Papers
    slide_data(prs)           # 4 - Pipeline
    slide_models(prs)         # 5 - 7 Models
    slide_features(prs)       # 6 - Features
    slide_demo_marker(prs,    # 7 - Demo Live
        "Demo: Live Dashboard",
        "PM2.5 now + multi-model predictions",
        "http://54.252.197.62:8502",
        "Station 56 | Metric cards | Chart")
    slide_demo_marker(prs,    # 8 - Demo Spike
        "Demo: Spike Replay",
        "Watch models predict a real spike hour-by-hour",
        "http://54.252.197.62:8502 -> Spike Replay",
        "Station 59 / 2025-01-24 | Auto-play | Scoreboard")
    slide_results(prs)        # 9 - Results (real data!)
    slide_when_to_use(prs)    # 10 - When to use
    slide_architecture(prs)   # 11 - Architecture
    slide_why_matters(prs)    # 12 - Impact
    slide_closing(prs)        # 13 - Closing
    slide_references(prs)     # 14 - References
    slide_qa(prs)             # 15 - Q&A

    out = Path("reports/FoonAlert_Presentation.pptx")
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"Saved: {out}")
    print(f"  {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
