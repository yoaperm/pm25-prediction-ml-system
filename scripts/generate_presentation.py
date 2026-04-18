"""
Generate PowerPoint presentation for PM2.5 Prediction ML System
FoonAlert — Term Project Presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette ──────────────────────────────────────────────────────
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)        # dark navy
BG_SECTION = RGBColor(0x16, 0x21, 0x3E)      # section header bg
ACCENT_BLUE = RGBColor(0x00, 0xD2, 0xFF)     # cyan accent
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x76)    # green accent
ACCENT_ORANGE = RGBColor(0xFF, 0x9F, 0x43)   # orange accent
ACCENT_RED = RGBColor(0xFF, 0x6B, 0x6B)      # red accent
ACCENT_PURPLE = RGBColor(0xA2, 0x9B, 0xFE)   # purple accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
MED_GRAY = RGBColor(0x88, 0x88, 0x99)

# modern clean palette
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT = RGBColor(0xF5, 0xF7, 0xFA)
PRIMARY = RGBColor(0x1E, 0x3A, 0x5F)          # dark blue
SECONDARY = RGBColor(0x3B, 0x82, 0xF6)        # bright blue
SUCCESS = RGBColor(0x10, 0xB9, 0x81)           # green
WARNING = RGBColor(0xF5, 0x9E, 0x0B)           # amber
DANGER = RGBColor(0xEF, 0x44, 0x44)            # red
TEXT_DARK = RGBColor(0x1F, 0x29, 0x37)
TEXT_MED = RGBColor(0x64, 0x74, 0x8B)
TEXT_LIGHT = RGBColor(0x94, 0xA3, 0xB8)
CARD_BG = RGBColor(0xF8, 0xFA, 0xFC)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_bg(slide, color=BG_WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width or Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, size=18, color=TEXT_DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return p


def add_bullet_text(tf, text, size=16, color=TEXT_DARK, level=0, bold=False, space_before=Pt(4), space_after=Pt(2)):
    p = tf.add_paragraph()
    p.level = level
    p.space_before = space_before
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Calibri"
    return p


def add_header_bar(slide):
    """Add a thin colored accent bar at the top of content slides"""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), SECONDARY)


def add_section_number(slide, number, left=Inches(0.7), top=Inches(0.5)):
    """Add a styled section number circle"""
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = SECONDARY
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(number)
    run.font.size = Pt(22)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = "Calibri"
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, PRIMARY)

# gradient-like overlay
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(2.5), RGBColor(0x15, 0x2F, 0x50))

# Title
tb = add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2))
set_text(tb.text_frame, "FoonAlert", size=52, color=WHITE, bold=True)
tb.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

tb2 = add_text_box(slide, Inches(1), Inches(2.6), Inches(11), Inches(1))
set_text(tb2.text_frame, "PM2.5 Air Quality Prediction System", size=32, color=ACCENT_BLUE, bold=False)

tb3 = add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8))
set_text(tb3.text_frame, "24-Hour Ahead Forecasting for 5 Bangkok Monitoring Stations", size=20, color=LIGHT_GRAY)

# Info cards at bottom
cards = [
    ("🤖", "5 ML Models", "Linear, Ridge, RF,\nXGBoost, LSTM"),
    ("📡", "5 Stations", "Real-time hourly data\nfrom AirBKK API"),
    ("⚡", "Triton Serving", "5-10ms inference\nONNX Runtime"),
    ("🔄", "Auto-Retrain", "PSI drift detection\n& monitoring"),
]
for i, (icon, title, desc) in enumerate(cards):
    x = Inches(1 + i * 2.9)
    card = add_rounded_rect(slide, x, Inches(4.8), Inches(2.6), Inches(2), RGBColor(0x24, 0x4A, 0x70))
    
    tb_icon = add_text_box(slide, x + Inches(0.15), Inches(4.95), Inches(2.3), Inches(0.5))
    set_text(tb_icon.text_frame, icon + "  " + title, size=16, color=WHITE, bold=True)
    
    tb_desc = add_text_box(slide, x + Inches(0.15), Inches(5.5), Inches(2.3), Inches(1))
    set_text(tb_desc.text_frame, desc, size=13, color=LIGHT_GRAY)

# Course info
tb_course = add_text_box(slide, Inches(1), Inches(7.0), Inches(11), Inches(0.4))
set_text(tb_course.text_frame, "2110555 Software Engineering for Machine Learning Systems  |  Academic Year 2568", size=12, color=MED_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2: AGENDA
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)

tb = add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8))
set_text(tb.text_frame, "Agenda", size=36, color=PRIMARY, bold=True)

agenda_items = [
    ("1", "Problem Statement & Use Case", "Real-world ML application & user interaction"),
    ("2", "Data & Experiment Management", "Data pipeline, feature engineering, MLflow tracking"),
    ("3", "Model Evaluation", "5 models comparison, metrics, ONNX export"),
    ("4", "System Architecture", "C4 model, ML system design, infrastructure"),
    ("5", "MLOps & Deployment", "CI/CD, Docker, Triton Inference Server"),
    ("6", "Monitoring & Drift Detection", "PSI, RMSE tracking, auto-retrain"),
    ("7", "Testing & Responsible ML", "Unit tests, debugging, fairness, limitations"),
    ("8", "Live Demo", "End-to-end demonstration"),
]

for i, (num, title, desc) in enumerate(agenda_items):
    y = Inches(1.3 + i * 0.72)
    
    # number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y + Inches(0.05), Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = SECONDARY if i % 2 == 0 else SUCCESS
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = num
    run.font.size = Pt(16)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = "Calibri"
    
    tb_title = add_text_box(slide, Inches(1.7), y, Inches(5), Inches(0.4))
    set_text(tb_title.text_frame, title, size=18, color=TEXT_DARK, bold=True)
    
    tb_desc = add_text_box(slide, Inches(1.7), y + Inches(0.32), Inches(8), Inches(0.35))
    set_text(tb_desc.text_frame, desc, size=14, color=TEXT_MED)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3: PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_section_number(slide, 1)

tb = add_text_box(slide, Inches(1.5), Inches(0.45), Inches(10), Inches(0.7))
set_text(tb.text_frame, "Problem Statement", size=32, color=PRIMARY, bold=True)

# Problem box
box = add_rounded_rect(slide, Inches(0.8), Inches(1.4), Inches(5.6), Inches(2.5), CARD_BG, border_color=RGBColor(0xE2, 0xE8, 0xF0))
tb = add_text_box(slide, Inches(1.0), Inches(1.5), Inches(5.2), Inches(2.3))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "The Challenge", size=20, color=DANGER, bold=True)
add_bullet_text(tf, "Bangkok faces severe PM2.5 pollution annually (Nov–Mar)", size=15, color=TEXT_DARK)
add_bullet_text(tf, "No accessible real-time forecasting system for the public", size=15, color=TEXT_DARK)
add_bullet_text(tf, "Existing monitoring only shows current values, not predictions", size=15, color=TEXT_DARK)

# Solution box
box2 = add_rounded_rect(slide, Inches(6.8), Inches(1.4), Inches(5.6), Inches(2.5), RGBColor(0xEC, 0xFD, 0xF5), border_color=SUCCESS)
tb = add_text_box(slide, Inches(7.0), Inches(1.5), Inches(5.2), Inches(2.3))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Our Solution", size=20, color=SUCCESS, bold=True)
add_bullet_text(tf, "Predict PM2.5 24 hours ahead for 5 Bangkok stations", size=15, color=TEXT_DARK)
add_bullet_text(tf, "Automated ML pipeline: ingest → train → serve → monitor", size=15, color=TEXT_DARK)
add_bullet_text(tf, "Auto-retrain when model performance degrades (drift)", size=15, color=TEXT_DARK)

# Foundation Model Data
tb_hdr = add_text_box(slide, Inches(0.8), Inches(4.2), Inches(5), Inches(0.5))
set_text(tb_hdr.text_frame, "Foundation Model — Data Splits", size=20, color=PRIMARY, bold=True)

splits = [
    ("Train", "2024 – mid 2025", "~1.5 years hourly data", SECONDARY),
    ("Validation", "Oct – Dec 2025", "3 months (hyperparameter tuning)", WARNING),
    ("Test", "Jan – Mar 2026", "3 months (final evaluation)", SUCCESS),
    ("Predict", "Last 24 hours", "Real-time inference input", ACCENT_PURPLE),
]
for i, (name, period, desc, color) in enumerate(splits):
    x = Inches(0.8 + i * 3.1)
    card = add_rounded_rect(slide, x, Inches(4.8), Inches(2.8), Inches(1.4), CARD_BG, border_color=color)
    
    # color bar on top of card
    add_rect(slide, x, Inches(4.8), Inches(2.8), Inches(0.06), color)
    
    tb = add_text_box(slide, x + Inches(0.15), Inches(4.95), Inches(2.5), Inches(0.4))
    set_text(tb.text_frame, name, size=16, color=color, bold=True)
    
    tb = add_text_box(slide, x + Inches(0.15), Inches(5.3), Inches(2.5), Inches(0.3))
    set_text(tb.text_frame, period, size=14, color=TEXT_DARK, bold=True)
    
    tb = add_text_box(slide, x + Inches(0.15), Inches(5.6), Inches(2.5), Inches(0.4))
    set_text(tb.text_frame, desc, size=12, color=TEXT_MED)

# Stations
tb_st = add_text_box(slide, Inches(0.8), Inches(6.4), Inches(12), Inches(0.6))
tf = tb_st.text_frame
set_text(tf, "Monitoring Stations:  ", size=14, color=TEXT_MED)
run = tf.paragraphs[0].add_run()
run.text = "Station 56  •  Station 57  •  Station 58  •  Station 59  •  Station 61   (Bangkok, Thailand)"
run.font.size = Pt(14)
run.font.color.rgb = PRIMARY
run.font.bold = True
run.font.name = "Calibri"


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4: USE CASE & USER INTERACTION
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_section_number(slide, 1)

tb = add_text_box(slide, Inches(1.5), Inches(0.45), Inches(10), Inches(0.7))
set_text(tb.text_frame, "Use Case & User Interaction", size=32, color=PRIMARY, bold=True)

# End User card
card1 = add_rounded_rect(slide, Inches(0.8), Inches(1.4), Inches(5.6), Inches(3.0), CARD_BG, border_color=SECONDARY)
add_rect(slide, Inches(0.8), Inches(1.4), Inches(5.6), Inches(0.06), SECONDARY)
tb = add_text_box(slide, Inches(1.0), Inches(1.55), Inches(5.2), Inches(2.7))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "👤  End User (Public)", size=20, color=SECONDARY, bold=True)
add_bullet_text(tf, "Access Streamlit Dashboard (port 8501)", size=15, color=TEXT_DARK)
add_bullet_text(tf, "Login with username/password (session-based auth)", size=15, color=TEXT_DARK)
add_bullet_text(tf, "View next-day PM2.5 forecast per station", size=15, color=TEXT_DARK)
add_bullet_text(tf, "See air quality level: Good / Moderate / Unhealthy", size=15, color=TEXT_DARK)
add_bullet_text(tf, "Upload CSV or manually enter 15+ days of history", size=15, color=TEXT_DARK)

# ML Engineer card
card2 = add_rounded_rect(slide, Inches(6.8), Inches(1.4), Inches(5.6), Inches(3.0), CARD_BG, border_color=SUCCESS)
add_rect(slide, Inches(6.8), Inches(1.4), Inches(5.6), Inches(0.06), SUCCESS)
tb = add_text_box(slide, Inches(7.0), Inches(1.55), Inches(5.2), Inches(2.7))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "🛠  ML Engineer / Admin", size=20, color=SUCCESS, bold=True)
add_bullet_text(tf, "Airflow UI (port 8080) — manage DAGs, trigger training", size=15, color=TEXT_DARK)
add_bullet_text(tf, "MLflow UI (port 5001) — view experiments, compare models", size=15, color=TEXT_DARK)
add_bullet_text(tf, "FastAPI (port 8001) — REST API for system integration", size=15, color=TEXT_DARK)
add_bullet_text(tf, "Monitor model performance & data quality", size=15, color=TEXT_DARK)

# Authentication Summary
tb_auth = add_text_box(slide, Inches(0.8), Inches(4.7), Inches(12), Inches(0.5))
set_text(tb_auth.text_frame, "Authentication & Security", size=20, color=PRIMARY, bold=True)

auth_items = [
    ("Streamlit", "Session-based login (VALID_USERS)", SECONDARY),
    ("FastAPI", "X-API-Key header (all endpoints except /health)", SUCCESS),
    ("Airflow", "Built-in auth (admin/admin)", WARNING),
]
for i, (svc, desc, color) in enumerate(auth_items):
    x = Inches(0.8 + i * 4.1)
    card = add_rounded_rect(slide, x, Inches(5.3), Inches(3.8), Inches(0.9), CARD_BG, border_color=color)
    add_rect(slide, x, Inches(5.3), Inches(0.08), Inches(0.9), color)
    tb = add_text_box(slide, x + Inches(0.25), Inches(5.35), Inches(3.4), Inches(0.35))
    set_text(tb.text_frame, svc, size=15, color=color, bold=True)
    tb = add_text_box(slide, x + Inches(0.25), Inches(5.65), Inches(3.4), Inches(0.35))
    set_text(tb.text_frame, desc, size=12, color=TEXT_MED)

# Live services
tb = add_text_box(slide, Inches(0.8), Inches(6.5), Inches(12), Inches(0.6))
tf = tb.text_frame
set_text(tf, "Live Production:  ", size=14, color=TEXT_MED)
run = tf.paragraphs[0].add_run()
run.text = "Streamlit :8501  |  FastAPI :8001  |  MLflow :5001  |  Airflow :8080  |  Triton :8010"
run.font.size = Pt(14)
run.font.color.rgb = PRIMARY
run.font.bold = True
run.font.name = "Calibri"


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5: DATA PIPELINE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_section_number(slide, 2)

tb = add_text_box(slide, Inches(1.5), Inches(0.45), Inches(10), Inches(0.7))
set_text(tb.text_frame, "Data Pipeline & Processing", size=32, color=PRIMARY, bold=True)

# Pipeline flow
steps = [
    ("AirBKK API", "Hourly PM2.5\ndata source", SECONDARY),
    ("Hourly Ingest\n(Airflow DAG)", "Fetch, validate,\nstore every hour", SUCCESS),
    ("PostgreSQL", "96K+ records\n5 stations", WARNING),
    ("Preprocessing", "ffill/bfill,\nclip [0, 500]", ACCENT_PURPLE),
    ("Feature Eng.", "19 features\nshift(1)", DANGER),
]
for i, (name, desc, color) in enumerate(steps):
    x = Inches(0.5 + i * 2.55)
    
    # Card
    card = add_rounded_rect(slide, x, Inches(1.4), Inches(2.2), Inches(1.8), CARD_BG, border_color=color)
    add_rect(slide, x, Inches(1.4), Inches(2.2), Inches(0.06), color)
    
    tb = add_text_box(slide, x + Inches(0.1), Inches(1.55), Inches(2.0), Inches(0.6))
    set_text(tb.text_frame, name, size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    
    tb = add_text_box(slide, x + Inches(0.1), Inches(2.15), Inches(2.0), Inches(0.8))
    set_text(tb.text_frame, desc, size=12, color=TEXT_MED, alignment=PP_ALIGN.CENTER)
    
    # Arrow between cards
    if i < len(steps) - 1:
        tb_arrow = add_text_box(slide, x + Inches(2.2), Inches(1.9), Inches(0.35), Inches(0.5))
        set_text(tb_arrow.text_frame, "→", size=24, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# Ingestion details
tb_hdr = add_text_box(slide, Inches(0.8), Inches(3.5), Inches(6), Inches(0.5))
set_text(tb_hdr.text_frame, "Hourly Ingestion Pipeline", size=20, color=PRIMARY, bold=True)

ingest_steps = [
    ("1. Fetch", "Call AirBKK API, convert Buddhist year (2569→2026), retry with backoff"),
    ("2. Validate", "Range checks: PM2.5 ∈ [0,500], RH ∈ [0,100], WS ≥ 0"),
    ("3. Store", "INSERT with UNIQUE(station_id, timestamp) — idempotent"),
    ("4. Monitor", "Check null rate (>50%), outliers (>10%), sensor drift (1h vs 7d baseline)"),
]
for i, (step, desc) in enumerate(ingest_steps):
    y = Inches(4.1 + i * 0.55)
    tb = add_text_box(slide, Inches(1.0), y, Inches(2), Inches(0.4))
    set_text(tb.text_frame, step, size=14, color=SECONDARY, bold=True)
    tb = add_text_box(slide, Inches(3.0), y, Inches(9), Inches(0.4))
    set_text(tb.text_frame, desc, size=13, color=TEXT_DARK)

# Data quality box
box = add_rounded_rect(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.8), RGBColor(0xFF, 0xFB, 0xEB), border_color=WARNING)
tb = add_text_box(slide, Inches(1.0), Inches(6.4), Inches(11.3), Inches(0.6))
tf = tb.text_frame
set_text(tf, "⚠  Data Quality:  ", size=14, color=WARNING, bold=True)
run = tf.paragraphs[0].add_run()
run.text = "Parameters collected: PM2.5, PM10, Temperature, Humidity, Wind Speed, Wind Direction  |  Stored in PostgreSQL with CHECK constraints"
run.font.size = Pt(13)
run.font.color.rgb = TEXT_DARK
run.font.name = "Calibri"


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6: FEATURE ENGINEERING
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_section_number(slide, 2)

tb = add_text_box(slide, Inches(1.5), Inches(0.45), Inches(10), Inches(0.7))
set_text(tb.text_frame, "Feature Engineering — 19 Features", size=32, color=PRIMARY, bold=True)

# Feature categories
categories = [
    ("Lag Features (6)", [
        "pm25_lag_1h", "pm25_lag_2h", "pm25_lag_3h",
        "pm25_lag_6h", "pm25_lag_12h", "pm25_lag_24h"
    ], SECONDARY),
    ("Rolling Statistics (6)", [
        "pm25_rolling_mean_6h", "pm25_rolling_mean_12h", "pm25_rolling_mean_24h",
        "pm25_rolling_std_6h", "pm25_rolling_std_12h", "pm25_rolling_std_24h"
    ], SUCCESS),
    ("Change Features (2)", [
        "pm25_diff_1h",
        "pm25_diff_24h",
    ], WARNING),
    ("Time Features (5)", [
        "hour", "day_of_week",
        "month", "day_of_year", "is_weekend"
    ], ACCENT_PURPLE),
]

for i, (cat_name, features, color) in enumerate(categories):
    x = Inches(0.5 + i * 3.15)
    
    card = add_rounded_rect(slide, x, Inches(1.4), Inches(2.9), Inches(3.5), CARD_BG, border_color=color)
    add_rect(slide, x, Inches(1.4), Inches(2.9), Inches(0.06), color)
    
    tb = add_text_box(slide, x + Inches(0.15), Inches(1.55), Inches(2.6), Inches(0.4))
    set_text(tb.text_frame, cat_name, size=15, color=color, bold=True)
    
    for j, feat in enumerate(features):
        tb = add_text_box(slide, x + Inches(0.15), Inches(2.05 + j * 0.38), Inches(2.6), Inches(0.35))
        set_text(tb.text_frame, "• " + feat, size=12, color=TEXT_DARK)

# Critical design decision
box = add_rounded_rect(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.8), RGBColor(0xFE, 0xF2, 0xF2), border_color=DANGER)
tb = add_text_box(slide, Inches(0.8), Inches(5.35), Inches(11.7), Inches(1.5))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "⚠  Critical Design: shift(1) on ALL features to prevent data leakage", size=18, color=DANGER, bold=True)
add_bullet_text(tf, "All lag and rolling features use shift(1) before computation", size=15, color=TEXT_DARK)
add_bullet_text(tf, "Ensures no information from the prediction hour is used as input", size=15, color=TEXT_DARK)
add_bullet_text(tf, "Code: df['pm25_rolling_mean_6h'] = df['pm25'].shift(1).rolling(6).mean()", size=14, color=TEXT_MED)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7: EXPERIMENT TRACKING
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_section_number(slide, 2)

tb = add_text_box(slide, Inches(1.5), Inches(0.45), Inches(10), Inches(0.7))
set_text(tb.text_frame, "Experiment Tracking & Environment", size=32, color=PRIMARY, bold=True)

# MLflow card
card = add_rounded_rect(slide, Inches(0.8), Inches(1.4), Inches(5.6), Inches(3.0), CARD_BG, border_color=SECONDARY)
add_rect(slide, Inches(0.8), Inches(1.4), Inches(5.6), Inches(0.06), SECONDARY)
tb = add_text_box(slide, Inches(1.0), Inches(1.55), Inches(5.2), Inches(2.7))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "📊  MLflow Tracking Server", size=20, color=SECONDARY, bold=True)
add_bullet_text(tf, "Backend: PostgreSQL (schema: mlflow)", size=15, color=TEXT_DARK)
add_bullet_text(tf, "Artifacts: /mlflow/artifacts (Docker volume)", size=15, color=TEXT_DARK)
add_bullet_text(tf, "Logs per training run:", size=15, color=TEXT_DARK)
add_bullet_text(tf, "Parameters — all hyperparameters per model", size=13, color=TEXT_MED, level=1)
add_bullet_text(tf, "Metrics — MAE, RMSE, R² on test set", size=13, color=TEXT_MED, level=1)
add_bullet_text(tf, "Artifacts — ONNX model files", size=13, color=TEXT_MED, level=1)
add_bullet_text(tf, "Tags — station_id, date range, best_model", size=13, color=TEXT_MED, level=1)

# Docker Compose stack
card2 = add_rounded_rect(slide, Inches(6.8), Inches(1.4), Inches(5.6), Inches(3.0), CARD_BG, border_color=SUCCESS)
add_rect(slide, Inches(6.8), Inches(1.4), Inches(5.6), Inches(0.06), SUCCESS)
tb = add_text_box(slide, Inches(7.0), Inches(1.55), Inches(5.2), Inches(2.7))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "🐳  Docker Compose (8 Services)", size=20, color=SUCCESS, bold=True)
services = [
    "PostgreSQL 15 — shared database",
    "MLflow — experiment tracking",
    "Airflow Webserver — DAG UI",
    "Airflow Scheduler — DAG execution",
    "Triton Server — ONNX inference",
    "FastAPI — prediction API",
    "Streamlit — web dashboard",
]
for svc in services:
    add_bullet_text(tf, svc, size=13, color=TEXT_DARK)

# Key env vars
tb_hdr = add_text_box(slide, Inches(0.8), Inches(4.7), Inches(12), Inches(0.5))
set_text(tb_hdr.text_frame, "Key Environment Variables", size=18, color=PRIMARY, bold=True)

env_vars = [
    ("INFERENCE_BACKEND", "triton / onnxruntime", "Inference engine selection"),
    ("OMP_NUM_THREADS=1", "MKL_NUM_THREADS=1", "Prevent XGBoost ↔ PyTorch thread conflict"),
    ("PYTORCH_DEVICE=cpu", "—", "Force CPU (Apple Silicon MPS crash fix)"),
    ("GRID_N_JOBS", "-1 (default)", "GridSearchCV parallelism"),
]

# Table header
headers = ["Variable", "Value", "Purpose"]
col_widths = [Inches(3.5), Inches(3), Inches(5.5)]
col_starts = [Inches(0.8), Inches(4.3), Inches(7.3)]
y_start = Inches(5.2)

add_rect(slide, Inches(0.8), y_start, Inches(11.7), Inches(0.4), PRIMARY)
for j, header in enumerate(headers):
    tb = add_text_box(slide, col_starts[j], y_start, col_widths[j], Inches(0.4))
    set_text(tb.text_frame, header, size=13, color=WHITE, bold=True)

for i, (var, val, purpose) in enumerate(env_vars):
    y = y_start + Inches(0.4 + i * 0.4)
    bg_color = CARD_BG if i % 2 == 0 else BG_WHITE
    add_rect(slide, Inches(0.8), y, Inches(11.7), Inches(0.4), bg_color)
    
    tb = add_text_box(slide, col_starts[0], y, col_widths[0], Inches(0.4))
    set_text(tb.text_frame, var, size=12, color=TEXT_DARK, bold=True)
    tb = add_text_box(slide, col_starts[1], y, col_widths[1], Inches(0.4))
    set_text(tb.text_frame, val, size=12, color=TEXT_MED)
    tb = add_text_box(slide, col_starts[2], y, col_widths[2], Inches(0.4))
    set_text(tb.text_frame, purpose, size=12, color=TEXT_DARK)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8: MODEL EVALUATION — 5 MODELS
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_section_number(slide, 3)

tb = add_text_box(slide, Inches(1.5), Inches(0.45), Inches(10), Inches(0.7))
set_text(tb.text_frame, "Model Evaluation — 5 Competing Models", size=32, color=PRIMARY, bold=True)

# Model cards
models = [
    ("Linear\nRegression", "Baseline", "No tuning\n(closed-form)", SECONDARY, "Simple,\ninterpretable"),
    ("Ridge\nRegression", "Regularized", "α ∈ [0.01..100]\nGridSearchCV", SUCCESS, "Handles correlated\nlag features"),
    ("Random\nForest", "Ensemble\n(Bagging)", "n_est, depth,\nGridSearchCV", WARNING, "Non-linear,\nfeature importance"),
    ("XGBoost", "Ensemble\n(Boosting)", "lr, depth, subsample\nGridSearchCV", DANGER, "SOTA tabular,\ngood for time-series"),
    ("LSTM", "Deep Learning\n(RNN)", "units, dropout, lr\nRandomizedSearch", ACCENT_PURPLE, "Temporal\ndependencies"),
]

for i, (name, mtype, tuning, color, why) in enumerate(models):
    x = Inches(0.3 + i * 2.55)
    
    card = add_rounded_rect(slide, x, Inches(1.4), Inches(2.3), Inches(4.0), CARD_BG, border_color=color)
    add_rect(slide, x, Inches(1.4), Inches(2.3), Inches(0.06), color)
    
    # Model name
    tb = add_text_box(slide, x + Inches(0.1), Inches(1.55), Inches(2.1), Inches(0.65))
    set_text(tb.text_frame, name, size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Type
    tb = add_text_box(slide, x + Inches(0.1), Inches(2.25), Inches(2.1), Inches(0.5))
    set_text(tb.text_frame, mtype, size=11, color=TEXT_MED, alignment=PP_ALIGN.CENTER)
    
    # Divider
    add_rect(slide, x + Inches(0.3), Inches(2.85), Inches(1.7), Inches(0.02), RGBColor(0xE2, 0xE8, 0xF0))
    
    # Tuning
    tb_lbl = add_text_box(slide, x + Inches(0.1), Inches(2.95), Inches(2.1), Inches(0.3))
    set_text(tb_lbl.text_frame, "Tuning", size=10, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)
    tb = add_text_box(slide, x + Inches(0.1), Inches(3.2), Inches(2.1), Inches(0.6))
    set_text(tb.text_frame, tuning, size=11, color=TEXT_DARK, alignment=PP_ALIGN.CENTER)
    
    # Why
    add_rect(slide, x + Inches(0.3), Inches(3.85), Inches(1.7), Inches(0.02), RGBColor(0xE2, 0xE8, 0xF0))
    tb_lbl = add_text_box(slide, x + Inches(0.1), Inches(3.95), Inches(2.1), Inches(0.3))
    set_text(tb_lbl.text_frame, "Why", size=10, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)
    tb = add_text_box(slide, x + Inches(0.1), Inches(4.2), Inches(2.1), Inches(0.6))
    set_text(tb.text_frame, why, size=11, color=TEXT_DARK, alignment=PP_ALIGN.CENTER)

# Strategy box
box = add_rounded_rect(slide, Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.4), CARD_BG, border_color=SECONDARY)
tb = add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.2))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Training Strategy", size=18, color=PRIMARY, bold=True)
add_bullet_text(tf, "CV: TimeSeriesSplit(n_splits=3) — preserves temporal order, prevents leakage", size=14, color=TEXT_DARK)
add_bullet_text(tf, "Scoring: neg_mean_absolute_error (MAE in µg/m³ — directly interpretable)", size=14, color=TEXT_DARK)
add_bullet_text(tf, "All 5 models train every cycle → best MAE on test set gets deployed (adaptive to data patterns)", size=14, color=TEXT_DARK)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9: METRICS & WHY ONNX
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_section_number(slide, 3)

tb = add_text_box(slide, Inches(1.5), Inches(0.45), Inches(10), Inches(0.7))
set_text(tb.text_frame, "Evaluation Metrics & Why ONNX", size=32, color=PRIMARY, bold=True)

# Metrics
card = add_rounded_rect(slide, Inches(0.8), Inches(1.4), Inches(5.6), Inches(2.8), CARD_BG, border_color=SECONDARY)
add_rect(slide, Inches(0.8), Inches(1.4), Inches(5.6), Inches(0.06), SECONDARY)
tb = add_text_box(slide, Inches(1.0), Inches(1.55), Inches(5.2), Inches(2.5))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "📏  Evaluation Metrics", size=20, color=SECONDARY, bold=True)
add_bullet_text(tf, "MAE (Primary) — Mean Absolute Error in µg/m³", size=15, color=TEXT_DARK, bold=True)
add_bullet_text(tf, "Average prediction error; directly interpretable", size=13, color=TEXT_MED, level=1)
add_bullet_text(tf, "RMSE — Root Mean Squared Error", size=15, color=TEXT_DARK, bold=True)
add_bullet_text(tf, "Penalizes large errors more heavily", size=13, color=TEXT_MED, level=1)
add_bullet_text(tf, "R² — Coefficient of Determination", size=15, color=TEXT_DARK, bold=True)
add_bullet_text(tf, "Proportion of variance explained by model", size=13, color=TEXT_MED, level=1)

# Why ONNX
card2 = add_rounded_rect(slide, Inches(6.8), Inches(1.4), Inches(5.6), Inches(2.8), CARD_BG, border_color=SUCCESS)
add_rect(slide, Inches(6.8), Inches(1.4), Inches(5.6), Inches(0.06), SUCCESS)
tb = add_text_box(slide, Inches(7.0), Inches(1.55), Inches(5.2), Inches(2.5))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "📦  Why ONNX-Only Deployment?", size=20, color=SUCCESS, bold=True)
add_bullet_text(tf, "Framework-agnostic — sklearn, XGBoost, PyTorch all export to ONNX", size=14, color=TEXT_DARK)
add_bullet_text(tf, "Triton native support — no extra conversion needed", size=14, color=TEXT_DARK)
add_bullet_text(tf, "No training deps at inference — lighter containers", size=14, color=TEXT_DARK)
add_bullet_text(tf, "Fast inference — optimized ONNX Runtime", size=14, color=TEXT_DARK)
add_bullet_text(tf, "Portable & versioned — file-based, easy rollback", size=14, color=TEXT_DARK)

# Export paths table
tb_hdr = add_text_box(slide, Inches(0.8), Inches(4.5), Inches(12), Inches(0.5))
set_text(tb_hdr.text_frame, "ONNX Export Paths", size=18, color=PRIMARY, bold=True)

headers = ["Model", "Export Library", "Function"]
col_starts = [Inches(0.8), Inches(4.5), Inches(8.2)]
col_widths = [Inches(3.7), Inches(3.7), Inches(4.3)]
y_start = Inches(5.0)

add_rect(slide, Inches(0.8), y_start, Inches(11.7), Inches(0.4), PRIMARY)
for j, h in enumerate(headers):
    tb = add_text_box(slide, col_starts[j], y_start, col_widths[j], Inches(0.4))
    set_text(tb.text_frame, h, size=13, color=WHITE, bold=True)

export_rows = [
    ("Linear, Ridge, Random Forest", "skl2onnx", "convert_sklearn()"),
    ("XGBoost", "onnxmltools", "convert_xgboost()"),
    ("LSTM (PyTorch)", "torch.onnx", "torch.onnx.export()"),
]
for i, (model, lib, func) in enumerate(export_rows):
    y = y_start + Inches(0.4 + i * 0.4)
    bg_color = CARD_BG if i % 2 == 0 else BG_WHITE
    add_rect(slide, Inches(0.8), y, Inches(11.7), Inches(0.4), bg_color)
    tb = add_text_box(slide, col_starts[0], y, col_widths[0], Inches(0.4))
    set_text(tb.text_frame, model, size=13, color=TEXT_DARK)
    tb = add_text_box(slide, col_starts[1], y, col_widths[1], Inches(0.4))
    set_text(tb.text_frame, lib, size=13, color=SECONDARY, bold=True)
    tb = add_text_box(slide, col_starts[2], y, col_widths[2], Inches(0.4))
    set_text(tb.text_frame, func, size=13, color=TEXT_MED)

# Deployment pointer
box = add_rounded_rect(slide, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.7), RGBColor(0xEC, 0xFD, 0xF5), border_color=SUCCESS)
tb = add_text_box(slide, Inches(1.0), Inches(6.45), Inches(11.3), Inches(0.6))
tf = tb.text_frame
set_text(tf, "Deployment Pointer:  active_model.json → ", size=14, color=SUCCESS, bold=True)
run = tf.paragraphs[0].add_run()
run.text = '{"onnx_file": "xgboost_2024-01-01_2025-12-31.onnx", "model_key": "xgboost", "is_lstm": false}'
run.font.size = Pt(12)
run.font.color.rgb = TEXT_MED
run.font.name = "Calibri"


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10: SYSTEM ARCHITECTURE — C4 CONTEXT
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_section_number(slide, 4)

tb = add_text_box(slide, Inches(1.5), Inches(0.45), Inches(10), Inches(0.7))
set_text(tb.text_frame, "System Architecture — C4 Context Diagram", size=32, color=PRIMARY, bold=True)

# Central system
sys_box = add_rounded_rect(slide, Inches(4.3), Inches(2.5), Inches(4.7), Inches(2.2), RGBColor(0xDB, 0xEA, 0xFE), border_color=SECONDARY)
tb = add_text_box(slide, Inches(4.5), Inches(2.6), Inches(4.3), Inches(0.5))
set_text(tb.text_frame, "PM2.5 Prediction System", size=18, color=SECONDARY, bold=True, alignment=PP_ALIGN.CENTER)
tb = add_text_box(slide, Inches(4.5), Inches(3.1), Inches(4.3), Inches(1.2))
set_text(tb.text_frame, "Predicts PM2.5 24h ahead\n5 Bangkok stations\nAuto-retrain on drift", size=14, color=TEXT_MED, alignment=PP_ALIGN.CENTER)

# Persons
for px, py, name, desc, color in [
    (Inches(0.5), Inches(1.0), "End User", "Views forecasts\nvia dashboard", SECONDARY),
    (Inches(0.5), Inches(4.2), "ML Engineer", "Manages pipeline\n& experiments", SUCCESS),
]:
    person = add_rounded_rect(slide, px, py, Inches(2.5), Inches(1.5), CARD_BG, border_color=color)
    tb = add_text_box(slide, px + Inches(0.1), py + Inches(0.1), Inches(2.3), Inches(0.4))
    set_text(tb.text_frame, "👤 " + name, size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    tb = add_text_box(slide, px + Inches(0.1), py + Inches(0.5), Inches(2.3), Inches(0.7))
    set_text(tb.text_frame, desc, size=12, color=TEXT_MED, alignment=PP_ALIGN.CENTER)

# External systems
for px, py, name, desc, color in [
    (Inches(10.0), Inches(1.0), "AirBKK API", "Hourly PM2.5\ndata source", WARNING),
    (Inches(10.0), Inches(3.0), "GitHub", "Source code\nCI/CD", TEXT_DARK),
    (Inches(10.0), Inches(5.0), "AWS EC2", "Production\nserver", DANGER),
]:
    ext = add_rounded_rect(slide, px, py, Inches(2.5), Inches(1.4), RGBColor(0xF1, 0xF5, 0xF9), border_color=color)
    tb = add_text_box(slide, px + Inches(0.1), py + Inches(0.1), Inches(2.3), Inches(0.4))
    set_text(tb.text_frame, name, size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    tb = add_text_box(slide, px + Inches(0.1), py + Inches(0.5), Inches(2.3), Inches(0.6))
    set_text(tb.text_frame, desc, size=12, color=TEXT_MED, alignment=PP_ALIGN.CENTER)

# Arrows as text
arrows = [
    (Inches(3.0), Inches(1.5), "→ HTTPS", TEXT_LIGHT),
    (Inches(3.0), Inches(4.7), "→ HTTPS", TEXT_LIGHT),
    (Inches(9.0), Inches(1.5), "← HTTP REST", TEXT_LIGHT),
    (Inches(9.0), Inches(3.5), "← SSH deploy", TEXT_LIGHT),
]
for ax, ay, label, c in arrows:
    tb = add_text_box(slide, ax, ay, Inches(1.5), Inches(0.3))
    set_text(tb.text_frame, label, size=11, color=c, alignment=PP_ALIGN.CENTER)

# Note
tb = add_text_box(slide, Inches(0.8), Inches(6.5), Inches(12), Inches(0.6))
set_text(tb.text_frame, "See full C4 diagrams (Level 1-4) in docs/C4_ARCHITECTURE.md with Mermaid source code", size=13, color=TEXT_LIGHT)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11: CONTAINER DIAGRAM
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_section_number(slide, 4)

tb = add_text_box(slide, Inches(1.5), Inches(0.45), Inches(10), Inches(0.7))
set_text(tb.text_frame, "System Architecture — Container Diagram", size=32, color=PRIMARY, bold=True)

# System boundary
add_rect(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.5), BG_WHITE, border_color=SECONDARY, border_width=Pt(2))
tb = add_text_box(slide, Inches(0.7), Inches(1.35), Inches(5), Inches(0.3))
set_text(tb.text_frame, "PM2.5 Prediction System (Docker Compose)", size=12, color=SECONDARY, bold=True)

containers = [
    # Row 1 — Presentation + API
    (Inches(0.8), Inches(1.8), "Streamlit\nDashboard", ":8501", "Python,\nStreamlit 1.41", RGBColor(0xFF, 0x4B, 0x4B)),
    (Inches(3.5), Inches(1.8), "FastAPI\nService", ":8001", "Python,\nFastAPI 0.115", RGBColor(0x00, 0x96, 0x88)),
    (Inches(6.2), Inches(1.8), "Triton\nServer", ":8010", "NVIDIA Triton,\nONNX Runtime", RGBColor(0x76, 0xB9, 0x00)),
    # Row 2 — Orchestration
    (Inches(0.8), Inches(3.7), "Airflow\nScheduler", ":8080", "Airflow 2.10\nDAG executor", RGBColor(0x01, 0x7C, 0xEE)),
    (Inches(3.5), Inches(3.7), "MLflow\nTracking", ":5001", "MLflow 2.16\nExperiment logs", RGBColor(0x01, 0x94, 0xE2)),
    (Inches(6.2), Inches(3.7), "PostgreSQL", ":5432", "PostgreSQL 15\n96K+ rows", RGBColor(0x33, 0x67, 0x91)),
    # Row 3 — Storage
    (Inches(9.0), Inches(1.8), "Model\nRepository", "File I/O", "ONNX files,\nactive_model.json", ACCENT_PURPLE),
    (Inches(9.0), Inches(3.7), "Triton\nModel Repo", "File I/O", "config.pbtxt,\nmodel.onnx", ACCENT_PURPLE),
]

for cx, cy, name, port, tech, color in containers:
    card = add_rounded_rect(slide, cx, cy, Inches(2.4), Inches(1.5), CARD_BG, border_color=color)
    add_rect(slide, cx, cy, Inches(2.4), Inches(0.05), color)
    
    tb = add_text_box(slide, cx + Inches(0.1), cy + Inches(0.1), Inches(2.2), Inches(0.5))
    set_text(tb.text_frame, name, size=12, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    
    tb = add_text_box(slide, cx + Inches(0.1), cy + Inches(0.65), Inches(2.2), Inches(0.3))
    set_text(tb.text_frame, port, size=10, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)
    
    tb = add_text_box(slide, cx + Inches(0.1), cy + Inches(0.9), Inches(2.2), Inches(0.5))
    set_text(tb.text_frame, tech, size=10, color=TEXT_MED, alignment=PP_ALIGN.CENTER)

# Key relationships (text arrows)
rels = [
    (Inches(2.7), Inches(2.4), "→", "HTTP"),
    (Inches(5.5), Inches(2.4), "→", "gRPC"),
    (Inches(8.2), Inches(2.4), "→", "File"),
    (Inches(2.7), Inches(4.3), "→", "HTTP"),
    (Inches(5.5), Inches(4.3), "→", "SQL"),
    (Inches(8.2), Inches(4.3), "→", "File"),
]
for rx, ry, arrow, label in rels:
    tb = add_text_box(slide, rx, ry, Inches(0.8), Inches(0.4))
    set_text(tb.text_frame, f"{arrow} {label}", size=9, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 12: ML PIPELINE COMPONENT DIAGRAM
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_section_number(slide, 4)

tb = add_text_box(slide, Inches(1.5), Inches(0.45), Inches(10), Inches(0.7))
set_text(tb.text_frame, "ML Pipeline — Component Diagram", size=32, color=PRIMARY, bold=True)

# Pipeline flow (horizontal)
components = [
    ("Ingest", "airbkk_client.py\nairflow_db.py", "Fetch hourly\nfrom AirBKK", SECONDARY),
    ("Preprocess", "preprocessing.py", "ffill, clip\n[0, 500]", SUCCESS),
    ("Features", "feature_eng.py", "19 features\nshift(1)", WARNING),
    ("Train", "train.py\nlstm_model.py", "5 models\nGridSearchCV", DANGER),
    ("Evaluate", "evaluate.py", "MAE, RMSE, R²\nvs production", ACCENT_PURPLE),
    ("Export", "export_onnx.py", "→ ONNX format", SECONDARY),
    ("Publish", "triton_utils.py", "→ Triton repo\nauto-reload", SUCCESS),
]

for i, (name, files, desc, color) in enumerate(components):
    x = Inches(0.2 + i * 1.85)
    
    card = add_rounded_rect(slide, x, Inches(1.4), Inches(1.65), Inches(2.5), CARD_BG, border_color=color)
    add_rect(slide, x, Inches(1.4), Inches(1.65), Inches(0.05), color)
    
    tb = add_text_box(slide, x + Inches(0.05), Inches(1.5), Inches(1.55), Inches(0.35))
    set_text(tb.text_frame, name, size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    
    tb = add_text_box(slide, x + Inches(0.05), Inches(1.9), Inches(1.55), Inches(0.5))
    set_text(tb.text_frame, files, size=9, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)
    
    tb = add_text_box(slide, x + Inches(0.05), Inches(2.5), Inches(1.55), Inches(0.7))
    set_text(tb.text_frame, desc, size=11, color=TEXT_DARK, alignment=PP_ALIGN.CENTER)
    
    if i < len(components) - 1:
        tb_arrow = add_text_box(slide, x + Inches(1.6), Inches(2.2), Inches(0.3), Inches(0.4))
        set_text(tb_arrow.text_frame, "→", size=18, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# Monitor feedback loop
monitor_box = add_rounded_rect(slide, Inches(3.5), Inches(4.3), Inches(6.3), Inches(1.6), RGBColor(0xFE, 0xF2, 0xF2), border_color=DANGER)
tb = add_text_box(slide, Inches(3.7), Inches(4.4), Inches(5.9), Inches(1.4))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "🔄  Drift Monitor (monitor.py + airflow_monitor.py)", size=16, color=DANGER, bold=True)
add_bullet_text(tf, "Daily check: rolling 14-day RMSE + PSI (10 bins)", size=14, color=TEXT_DARK)
add_bullet_text(tf, "If RMSE > 13.0 OR PSI > 0.2 → trigger retrain pipeline", size=14, color=TEXT_DARK)
add_bullet_text(tf, "New model deployed only if MAE improves over production", size=14, color=TEXT_DARK)

# Feedback arrow
tb = add_text_box(slide, Inches(1.5), Inches(4.6), Inches(2), Inches(0.5))
set_text(tb.text_frame, "↻ Retrain trigger", size=13, color=DANGER, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 13: DEPLOYMENT & MLOPS
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_section_number(slide, 5)

tb = add_text_box(slide, Inches(1.5), Inches(0.45), Inches(10), Inches(0.7))
set_text(tb.text_frame, "MLOps & Deployment", size=32, color=PRIMARY, bold=True)

# Triton card
card = add_rounded_rect(slide, Inches(0.8), Inches(1.4), Inches(5.6), Inches(2.5), CARD_BG, border_color=SUCCESS)
add_rect(slide, Inches(0.8), Inches(1.4), Inches(5.6), Inches(0.06), SUCCESS)
tb = add_text_box(slide, Inches(1.0), Inches(1.55), Inches(5.2), Inches(2.2))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "⚡  Triton Inference Server", size=20, color=SUCCESS, bold=True)
add_bullet_text(tf, "NVIDIA Triton 24.08 with ONNX Runtime backend", size=14, color=TEXT_DARK)
add_bullet_text(tf, "5-10ms inference latency per request", size=14, color=TEXT_DARK)
add_bullet_text(tf, "Dynamic batching (max_batch_size=32)", size=14, color=TEXT_DARK)
add_bullet_text(tf, "Auto-reload: polls model repo every 30 seconds", size=14, color=TEXT_DARK)
add_bullet_text(tf, "Zero-downtime deploy: publish ONNX → auto-picked up", size=14, color=TEXT_DARK)
add_bullet_text(tf, "Fallback: FastAPI uses onnxruntime directly", size=14, color=TEXT_DARK)

# CI/CD card
card2 = add_rounded_rect(slide, Inches(6.8), Inches(1.4), Inches(5.6), Inches(2.5), CARD_BG, border_color=SECONDARY)
add_rect(slide, Inches(6.8), Inches(1.4), Inches(5.6), Inches(0.06), SECONDARY)
tb = add_text_box(slide, Inches(7.0), Inches(1.55), Inches(5.2), Inches(2.2))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "🔄  CI/CD Pipeline (GitHub Actions)", size=20, color=SECONDARY, bold=True)
add_bullet_text(tf, "CI — Test & Lint: pytest + ruff", size=14, color=TEXT_DARK)
add_bullet_text(tf, "Triggers: push main/develop, PRs to main", size=14, color=TEXT_DARK)
add_bullet_text(tf, "CD — Deploy to EC2:", size=14, color=TEXT_DARK)
add_bullet_text(tf, "SSH to EC2 → git pull → docker compose up --build", size=13, color=TEXT_MED, level=1)
add_bullet_text(tf, "Auto-deploy on push to main branch", size=13, color=TEXT_MED, level=1)

# Airflow DAGs
tb_hdr = add_text_box(slide, Inches(0.8), Inches(4.2), Inches(12), Inches(0.5))
set_text(tb_hdr.text_frame, "Airflow DAGs — Orchestration", size=20, color=PRIMARY, bold=True)

dags = [
    ("pm25_hourly_ingest", "Every hour (0 * * * *)", "Fetch from AirBKK, validate, store", SECONDARY),
    ("pm25_24h_training", "Manual trigger", "Train 5 models, deploy best per station", SUCCESS),
    ("pm25_24h_pipeline", "Daily 01:00 UTC", "Monitor drift → trigger retrain if needed", WARNING),
    ("pm25_api_prediction", "Manual/Scheduled", "Predict via API for date range", ACCENT_PURPLE),
]

headers = ["DAG", "Schedule", "Purpose"]
col_starts = [Inches(0.8), Inches(4.5), Inches(7.5)]
col_widths = [Inches(3.7), Inches(3.0), Inches(5.0)]
y_start = Inches(4.8)

add_rect(slide, Inches(0.8), y_start, Inches(11.7), Inches(0.4), PRIMARY)
for j, h in enumerate(headers):
    tb = add_text_box(slide, col_starts[j], y_start, col_widths[j], Inches(0.4))
    set_text(tb.text_frame, h, size=13, color=WHITE, bold=True)

for i, (dag, sched, purpose, color) in enumerate(dags):
    y = y_start + Inches(0.4 + i * 0.45)
    bg_color = CARD_BG if i % 2 == 0 else BG_WHITE
    add_rect(slide, Inches(0.8), y, Inches(11.7), Inches(0.45), bg_color)
    
    tb = add_text_box(slide, col_starts[0], y, col_widths[0], Inches(0.45))
    set_text(tb.text_frame, dag, size=12, color=color, bold=True)
    tb = add_text_box(slide, col_starts[1], y, col_widths[1], Inches(0.45))
    set_text(tb.text_frame, sched, size=12, color=TEXT_MED)
    tb = add_text_box(slide, col_starts[2], y, col_widths[2], Inches(0.45))
    set_text(tb.text_frame, purpose, size=12, color=TEXT_DARK)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 14: MONITORING & DRIFT
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_section_number(slide, 6)

tb = add_text_box(slide, Inches(1.5), Inches(0.45), Inches(10), Inches(0.7))
set_text(tb.text_frame, "Monitoring & Drift Detection", size=32, color=PRIMARY, bold=True)

# Performance Monitoring
card = add_rounded_rect(slide, Inches(0.8), Inches(1.4), Inches(5.6), Inches(2.8), CARD_BG, border_color=DANGER)
add_rect(slide, Inches(0.8), Inches(1.4), Inches(5.6), Inches(0.06), DANGER)
tb = add_text_box(slide, Inches(1.0), Inches(1.55), Inches(5.2), Inches(2.5))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "📊  Model Performance Monitoring", size=20, color=DANGER, bold=True)
add_bullet_text(tf, "RMSE — rolling 14-day window", size=15, color=TEXT_DARK, bold=True)
add_bullet_text(tf, "Threshold: > 13.0 µg/m³ → trigger retrain", size=14, color=TEXT_MED, level=1)
add_bullet_text(tf, "PSI — Population Stability Index", size=15, color=TEXT_DARK, bold=True)
add_bullet_text(tf, "Measures distribution shift (predicted vs actual)", size=14, color=TEXT_MED, level=1)
add_bullet_text(tf, "10 percentile-based bins", size=14, color=TEXT_MED, level=1)
add_bullet_text(tf, "< 0.1 = Stable  |  0.1–0.2 = Monitor  |  > 0.2 = Retrain", size=14, color=TEXT_MED, level=1)

# Data Quality
card2 = add_rounded_rect(slide, Inches(6.8), Inches(1.4), Inches(5.6), Inches(2.8), CARD_BG, border_color=WARNING)
add_rect(slide, Inches(6.8), Inches(1.4), Inches(5.6), Inches(0.06), WARNING)
tb = add_text_box(slide, Inches(7.0), Inches(1.55), Inches(5.2), Inches(2.5))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "⚠  Data Quality Monitoring", size=20, color=WARNING, bold=True)
add_bullet_text(tf, "Runs after every hourly ingestion", size=15, color=TEXT_DARK)
add_bullet_text(tf, "Null rate check — alert if > 50%", size=14, color=TEXT_DARK)
add_bullet_text(tf, "Outlier rate check — alert if > 10%", size=14, color=TEXT_DARK)
add_bullet_text(tf, "Extreme values — alert if mean > 400", size=14, color=TEXT_DARK)
add_bullet_text(tf, "Sensor drift — compare 1h vs 7-day baseline", size=14, color=TEXT_DARK)
add_bullet_text(tf, "> 25% drift = MODERATE, > 50% = SEVERE", size=14, color=TEXT_MED, level=1)

# Auto-retrain flow
tb_hdr = add_text_box(slide, Inches(0.8), Inches(4.5), Inches(12), Inches(0.5))
set_text(tb_hdr.text_frame, "Auto-Retrain Flow", size=20, color=PRIMARY, bold=True)

flow_steps = [
    ("Daily Check\n01:00 UTC", "Query 14-day\npredictions\n& actuals", "Compute\nRMSE + PSI", "RMSE>13\nOR PSI>0.2?", "Train 5\nmodels", "Best MAE\nimproves?", "Deploy\nto Triton"),
]
colors = [SECONDARY, TEXT_DARK, TEXT_DARK, WARNING, DANGER, WARNING, SUCCESS]
labels = ["Daily Check\n01:00 UTC", "Query 14-day\npredictions", "Compute\nRMSE + PSI", "Threshold\ncheck", "Train 5\nmodels", "Compare vs\nproduction", "Deploy new\nmodel"]

for i, (label, color) in enumerate(zip(labels, colors)):
    x = Inches(0.4 + i * 1.8)
    card = add_rounded_rect(slide, x, Inches(5.1), Inches(1.6), Inches(1.2), CARD_BG, border_color=color)
    tb = add_text_box(slide, x + Inches(0.05), Inches(5.15), Inches(1.5), Inches(1.0))
    set_text(tb.text_frame, label, size=11, color=color, bold=(i==0 or i==6), alignment=PP_ALIGN.CENTER)
    
    if i < len(labels) - 1:
        tb = add_text_box(slide, x + Inches(1.55), Inches(5.4), Inches(0.3), Inches(0.4))
        set_text(tb.text_frame, "→", size=16, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 15: TESTING & RESPONSIBLE ML
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_section_number(slide, 7)

tb = add_text_box(slide, Inches(1.5), Inches(0.45), Inches(10), Inches(0.7))
set_text(tb.text_frame, "Testing, Maintainability & Responsible ML", size=32, color=PRIMARY, bold=True)

# Testing
card = add_rounded_rect(slide, Inches(0.8), Inches(1.4), Inches(3.7), Inches(2.8), CARD_BG, border_color=SECONDARY)
add_rect(slide, Inches(0.8), Inches(1.4), Inches(3.7), Inches(0.06), SECONDARY)
tb = add_text_box(slide, Inches(1.0), Inches(1.55), Inches(3.3), Inches(2.5))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "🧪  Testing", size=18, color=SECONDARY, bold=True)
add_bullet_text(tf, "Unit tests (pytest)", size=14, color=TEXT_DARK, bold=True)
add_bullet_text(tf, "test_handle_missing_ffill", size=12, color=TEXT_MED, level=1)
add_bullet_text(tf, "test_handle_missing_interpolate", size=12, color=TEXT_MED, level=1)
add_bullet_text(tf, "test_remove_outliers", size=12, color=TEXT_MED, level=1)
add_bullet_text(tf, "Linting: ruff check", size=14, color=TEXT_DARK, bold=True)
add_bullet_text(tf, "CI: GitHub Actions", size=14, color=TEXT_DARK, bold=True)

# Debugging
card2 = add_rounded_rect(slide, Inches(4.8), Inches(1.4), Inches(3.7), Inches(2.8), CARD_BG, border_color=SUCCESS)
add_rect(slide, Inches(4.8), Inches(1.4), Inches(3.7), Inches(0.06), SUCCESS)
tb = add_text_box(slide, Inches(5.0), Inches(1.55), Inches(3.3), Inches(2.5))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "🔍  Debugging", size=18, color=SUCCESS, bold=True)
add_bullet_text(tf, "Airflow UI — task logs, XCom", size=13, color=TEXT_DARK)
add_bullet_text(tf, "MLflow UI — experiment comparison", size=13, color=TEXT_DARK)
add_bullet_text(tf, "CSV logs — predictions, actuals, monitoring", size=13, color=TEXT_DARK)
add_bullet_text(tf, "Docker logs — per service", size=13, color=TEXT_DARK)
add_bullet_text(tf, "Health endpoints — /health", size=13, color=TEXT_DARK)

# Responsible ML
card3 = add_rounded_rect(slide, Inches(8.8), Inches(1.4), Inches(3.7), Inches(2.8), CARD_BG, border_color=ACCENT_PURPLE)
add_rect(slide, Inches(8.8), Inches(1.4), Inches(3.7), Inches(0.06), ACCENT_PURPLE)
tb = add_text_box(slide, Inches(9.0), Inches(1.55), Inches(3.3), Inches(2.5))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "⚖  Responsible ML", size=18, color=ACCENT_PURPLE, bold=True)
add_bullet_text(tf, "No personal data collected", size=13, color=TEXT_DARK)
add_bullet_text(tf, "Per-station models (reduce geographic bias)", size=13, color=TEXT_DARK)
add_bullet_text(tf, "Feature importance for explainability", size=13, color=TEXT_DARK)
add_bullet_text(tf, "API key authentication", size=13, color=TEXT_DARK)
add_bullet_text(tf, "Clear system limitations documented", size=13, color=TEXT_DARK)

# Technical Debt
tb_hdr = add_text_box(slide, Inches(0.8), Inches(4.5), Inches(12), Inches(0.5))
set_text(tb_hdr.text_frame, "Technical Debt & Improvements", size=18, color=PRIMARY, bold=True)

debts = [
    ("Tests", "Preprocessing only", "Add feature eng, API, model eval tests"),
    ("Data Versioning", "No DVC", "Add DVC for training data tracking"),
    ("Model Registry", "File-based (active_model.json)", "Migrate to MLflow Model Registry"),
    ("Scaling", "Docker Compose (single node)", "Kubernetes for horizontal scaling"),
]

headers = ["Area", "Current", "Improvement"]
col_starts = [Inches(0.8), Inches(3.5), Inches(7.5)]
col_widths = [Inches(2.7), Inches(4.0), Inches(5.0)]
y_start = Inches(5.0)

add_rect(slide, Inches(0.8), y_start, Inches(11.7), Inches(0.4), PRIMARY)
for j, h in enumerate(headers):
    tb = add_text_box(slide, col_starts[j], y_start, col_widths[j], Inches(0.4))
    set_text(tb.text_frame, h, size=13, color=WHITE, bold=True)

for i, (area, current, improvement) in enumerate(debts):
    y = y_start + Inches(0.4 + i * 0.4)
    bg_color = CARD_BG if i % 2 == 0 else BG_WHITE
    add_rect(slide, Inches(0.8), y, Inches(11.7), Inches(0.4), bg_color)
    tb = add_text_box(slide, col_starts[0], y, col_widths[0], Inches(0.4))
    set_text(tb.text_frame, area, size=12, color=TEXT_DARK, bold=True)
    tb = add_text_box(slide, col_starts[1], y, col_widths[1], Inches(0.4))
    set_text(tb.text_frame, current, size=12, color=DANGER)
    tb = add_text_box(slide, col_starts[2], y, col_widths[2], Inches(0.4))
    set_text(tb.text_frame, improvement, size=12, color=SUCCESS)

# Limitations
box = add_rounded_rect(slide, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5), RGBColor(0xFF, 0xFB, 0xEB), border_color=WARNING)
tb = add_text_box(slide, Inches(1.0), Inches(6.75), Inches(11.3), Inches(0.4))
set_text(tb.text_frame, "Limitations: Single data source (AirBKK) • 5 stations only • 24h horizon • Weak on extreme events (>200 µg/m³)", size=13, color=TEXT_DARK)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 16: TRADE-OFFS
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_section_number(slide, 4)

tb = add_text_box(slide, Inches(1.5), Inches(0.45), Inches(10), Inches(0.7))
set_text(tb.text_frame, "Design Decisions & Trade-offs", size=32, color=PRIMARY, bold=True)

tradeoffs = [
    ("ONNX-only inference", "Framework-agnostic, fast, portable", "Must export each model type differently"),
    ("Triton Inference Server", "Low latency (5-10ms), auto-batching", "Adds operational complexity"),
    ("5 competing models", "Adaptive — picks best for current data", "Training takes 15-20h per 5 stations"),
    ("Shared PostgreSQL", "Reduces ops overhead, single backup", "Single point of failure"),
    ("shift(1) on all features", "100% leakage prevention", "Loses 1 row per feature"),
    ("PSI + RMSE monitoring", "Catches both drift and accuracy drop", "Needs matched actual data"),
    ("Docker Compose deploy", "Simple, reproducible, one command", "No horizontal scaling (vs K8s)"),
    ("File-based model versioning", "Simple, Triton polls directly", "No registry UI, manual rollback"),
]

headers = ["Decision", "Benefit ✅", "Trade-off ⚠️"]
col_starts = [Inches(0.8), Inches(4.3), Inches(8.8)]
col_widths = [Inches(3.5), Inches(4.5), Inches(4.0)]
y_start = Inches(1.3)

add_rect(slide, Inches(0.8), y_start, Inches(11.7), Inches(0.45), PRIMARY)
for j, h in enumerate(headers):
    tb = add_text_box(slide, col_starts[j], y_start, col_widths[j], Inches(0.45))
    set_text(tb.text_frame, h, size=14, color=WHITE, bold=True)

for i, (decision, benefit, tradeoff) in enumerate(tradeoffs):
    y = y_start + Inches(0.45 + i * 0.65)
    bg_color = CARD_BG if i % 2 == 0 else BG_WHITE
    add_rect(slide, Inches(0.8), y, Inches(11.7), Inches(0.65), bg_color)
    
    tb = add_text_box(slide, col_starts[0], y + Inches(0.05), col_widths[0], Inches(0.55))
    set_text(tb.text_frame, decision, size=13, color=PRIMARY, bold=True)
    tb = add_text_box(slide, col_starts[1], y + Inches(0.05), col_widths[1], Inches(0.55))
    set_text(tb.text_frame, benefit, size=12, color=SUCCESS)
    tb = add_text_box(slide, col_starts[2], y + Inches(0.05), col_widths[2], Inches(0.55))
    set_text(tb.text_frame, tradeoff, size=12, color=WARNING)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 17: LIVE DEMO
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)
add_section_number(slide, 8, left=Inches(5.9), top=Inches(1.5))

tb = add_text_box(slide, Inches(2), Inches(2.2), Inches(9), Inches(1))
set_text(tb.text_frame, "Live Demo", size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

tb = add_text_box(slide, Inches(2), Inches(3.2), Inches(9), Inches(0.6))
set_text(tb.text_frame, "End-to-end demonstration of the PM2.5 Prediction System", size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Demo items
demo_items = [
    ("🖥", "Streamlit — Login → Predict → Results → Monitoring"),
    ("🔄", "Airflow — Ingest DAG running, trigger training"),
    ("📊", "MLflow — Experiment tracking, model comparison"),
    ("⚡", "API — curl /predict with API key → response"),
    ("🔧", "CI/CD — GitHub Actions deploy flow"),
]

for i, (icon, text) in enumerate(demo_items):
    y = Inches(4.1 + i * 0.55)
    tb = add_text_box(slide, Inches(3), y, Inches(7), Inches(0.5))
    tf = tb.text_frame
    set_text(tf, f"{icon}  {text}", size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 18: THANK YOU
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)

tb = add_text_box(slide, Inches(2), Inches(2.0), Inches(9), Inches(1))
set_text(tb.text_frame, "Thank You", size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

tb = add_text_box(slide, Inches(2), Inches(3.2), Inches(9), Inches(0.6))
set_text(tb.text_frame, "FoonAlert — PM2.5 Prediction ML System", size=24, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

tb = add_text_box(slide, Inches(2), Inches(4.2), Inches(9), Inches(0.5))
set_text(tb.text_frame, "Questions?", size=28, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Service URLs
urls = [
    ("Dashboard", "http://43.209.207.187:8501"),
    ("API", "http://43.209.207.187:8001"),
    ("MLflow", "http://43.209.207.187:5001"),
    ("Airflow", "http://43.209.207.187:8080"),
]
for i, (name, url) in enumerate(urls):
    x = Inches(1.5 + i * 2.8)
    card = add_rounded_rect(slide, x, Inches(5.3), Inches(2.5), Inches(1.0), RGBColor(0x24, 0x4A, 0x70))
    tb = add_text_box(slide, x + Inches(0.1), Inches(5.35), Inches(2.3), Inches(0.35))
    set_text(tb.text_frame, name, size=14, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    tb = add_text_box(slide, x + Inches(0.1), Inches(5.7), Inches(2.3), Inches(0.4))
    set_text(tb.text_frame, url, size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════
output_path = "docs/FoonAlert_PM25_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to {output_path}")
print(f"   Total slides: {len(prs.slides)}")
